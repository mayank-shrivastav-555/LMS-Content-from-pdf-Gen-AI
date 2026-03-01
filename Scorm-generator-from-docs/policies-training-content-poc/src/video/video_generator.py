from pptx import Presentation
import os

HEADING_COLOR = (103, 255, 242)
HEADING_HIGHLIGHT_COLOR = (103, 255, 242)
BULLET_COLOR = (255, 255, 255)
BG_COLOR = (240, 240, 240)

class VideoGenerator:
    def ppt_to_images(self, ppt_path: str, output_dir: str, gemini_model=None, icon_paths=None, default_img_path=None, use_image_background=False, intro_img_path=None) -> list:
        """
        Converts PPT slides to images using python-pptx and Pillow, overlays relevant images fetched via SlideImageFetcher.
        Returns list of image file paths.
        Enhancements:
        - Prevents text overlap with images
        - Optionally uses image as background (toggle)
        - Bullet headings (e.g. 'OUR PURPOSE:') split after colon
        - Proper font sizes, indentation, wrapping
        """
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        import requests
        from utils.slide_image_fetcher import SlideImageFetcher
        import os
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        prs = Presentation(ppt_path)
        image_files = []
        width = int(prs.slide_width / 914400 * 96)
        height = int(prs.slide_height / 914400 * 96)
        if width <= 0 or width > 1920:
            width = 1280
        if height <= 0 or height > 1080:
            height = 720
        try:
            heading_font = ImageFont.truetype("arialbd.ttf", 54)
            bullet_font = ImageFont.truetype("arial.ttf", 32)
        except:
            heading_font = ImageFont.load_default()
            bullet_font = ImageFont.load_default()
        slides = list(prs.slides)
        fetcher = SlideImageFetcher()
        if icon_paths is None:
            icon_dir = os.path.join(output_dir, 'icons')
            if os.path.exists(icon_dir):
                icon_paths = [os.path.join(icon_dir, f) for f in os.listdir(icon_dir) if f.endswith('.png')]
            else:
                icon_paths = []
        margin_x = 60  # Improved left margin
        margin_y = 40
        bullet_indent = 40  # Slightly reduced indent for more space
        bullet_spacing = 44
        for idx, slide in enumerate(slides):
            is_thankyou_slide = (idx == len(slides) - 1)
            if is_thankyou_slide:
                title_text = "Thank You!"
            else:
                title_text = slide.shapes.title.text if slide.shapes.title and slide.shapes.title.text else ""
                if not title_text:
                    # Use first bullet as heading if no explicit title
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
                            title_text = shape.text.split('\n')[0].strip()
                            break
            # --- Image background toggle ---
            img = None
            print(f"Processing slide{idx}, {idx+1}/{len(slides)}: '{title_text}'")
            # For each slide
            if idx == 0:
                # Keep the intro slide image as-is
                try:
                    if os.path.exists(intro_img_path):
                        intro_img = Image.open(intro_img_path).convert("RGB")
                        intro_img = intro_img.resize((width, height))
                        img = intro_img.copy()
                    else:
                        raise FileNotFoundError(f"Intro image not found at {intro_img_path}")
                except Exception as e:
                    print(f"Error loading intro image: {e}")
                    img = Image.new('RGB', (width, height), color=BG_COLOR)  # fallback plain image
            else:
                if use_image_background and default_img_path and os.path.exists(default_img_path):
                    try:
                        default_img = Image.open(default_img_path).convert("RGBA")
                        default_img = default_img.resize((width, height))
                        img = default_img.copy()
                    except Exception as e:
                        print(f"Default image paste error for slide {idx}: {e}")
                        img = Image.new('RGB', (width, height), color=BG_COLOR)
                    overlay = Image.new('RGBA', (width, height), (38, 52, 126, 100))
                    img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
                else:
                    img = Image.new('RGB', (width, height), color=BG_COLOR)

            # Wrap heading to avoid cutoff
            def wrap_heading(text, font, max_width, draw):
                import textwrap
                lines = []
                words = text.split()
                current_line = ""
                for word in words:
                    test_line = current_line + (" " if current_line else "") + word
                    try:
                        bbox = draw.textbbox((0, 0), test_line, font=font)
                        w = bbox[2] - bbox[0]
                    except Exception:
                        w = len(test_line) * 20
                    if w > max_width and current_line:
                        lines.append(current_line)
                        current_line = word
                    else:
                        current_line = test_line
                if current_line:
                    lines.append(current_line)
                return lines
            heading_margin_x = margin_x
            heading_max_width = width - 2 * heading_margin_x
            draw = ImageDraw.Draw(img)
            heading_lines = wrap_heading(title_text, heading_font, heading_max_width, draw)
            # Set max_text_width to allow 60px right margin
            max_text_width = width - (margin_x + bullet_indent + 60)
            y = margin_y
            # Title
            if title_text:
                if is_thankyou_slide:
                    # Center Thank You text both vertically and horizontally, only blue
                    bbox = draw.textbbox((0, 0), title_text, font=heading_font)
                    th = bbox[3] - bbox[1]
                    tw = bbox[2] - bbox[0]
                    x = int((width - tw) / 2)
                    y = int((height - th) / 2)
                    draw.text((x, y), title_text, font=heading_font, fill=HEADING_COLOR)
                else:
                    # Center heading vertically in top 1/3 of slide
                    total_heading_height = 0
                    for line in heading_lines:
                        try:
                            bbox = draw.textbbox((heading_margin_x, 0), line, font=heading_font)
                            th = bbox[3] - bbox[1]
                        except Exception:
                            th = 60
                        total_heading_height += th + 8
                    heading_y = margin_y + int((height * 0.33 - total_heading_height) / 2)
                    for line in heading_lines:
                        draw.text((heading_margin_x, heading_y), line, font=heading_font, fill=HEADING_COLOR)
                        try:
                            bbox = draw.textbbox((heading_margin_x, heading_y), line, font=heading_font)
                            th = bbox[3] - bbox[1]
                        except Exception:
                            th = 60
                        heading_y += th + 8
                    y = heading_y + 32  # Add extra vertical spacing after heading before bullets
            # Overlay corporate icon (bottom left, 80x80)
            # if icon_paths and len(icon_paths) > 0:
            #     try:
            #         icon_path = icon_paths[idx % len(icon_paths)]
            #         icon_img = Image.open(icon_path).convert("RGBA").resize((80, 80))
            #         img.paste(icon_img, (margin_x, height - 120), icon_img)
            #     except Exception as e:
            #         print(f"Icon overlay error for slide {idx}: {e}")
            # Bullets (wrap text, avoid image area)
            if not is_thankyou_slide:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
                        for line in shape.text.split('\n'):
                            # Enhancement: If line contains heading (e.g., 'OUR PURPOSE:'), split and put content on next line
                            if ':' in line:
                                parts = line.split(':', 1)
                                heading = parts[0].strip()
                                content = parts[1].strip()
                                draw.text((margin_x + bullet_indent, y), heading + ':', font=bullet_font, fill=HEADING_COLOR)
                                try:
                                    bbox = draw.textbbox((margin_x + bullet_indent, y), heading + ':', font=bullet_font)
                                    h = bbox[3] - bbox[1]
                                except Exception:
                                    h = bullet_spacing
                                y += h + 4
                                line = content
                            words = line.split()
                            current_line = ""
                            for word in words:
                                test_line = current_line + (" " if current_line else "") + word
                                try:
                                    bbox = draw.textbbox((margin_x + bullet_indent, y), test_line, font=bullet_font)
                                    w = bbox[2] - bbox[0]
                                except Exception:
                                    w = len(test_line) * 20
                                if w > max_text_width:
                                    draw.text((margin_x + bullet_indent, y), current_line, font=bullet_font, fill=BULLET_COLOR)
                                    y += bullet_spacing
                                    current_line = word
                                else:
                                    current_line = test_line
                            if current_line:
                                draw.text((margin_x + bullet_indent, y), current_line, font=bullet_font, fill=BULLET_COLOR)
                                y += bullet_spacing
            img_path = os.path.join(output_dir, f'slide_{idx}.png')
            img.save(img_path)
            image_files.append(img_path)
        return image_files


    def create_video(self, image_files: list, audio_files: list, output_path: str, overlay_texts: list = None, background_music_path: str = None, slide_pause_duration: int = 1):
        """
        Creates a video by combining each slide image with its corresponding audio segment, then concatenates all segments with fade transitions and animated overlays.
        Uses ffmpeg complex filter for batching and transitions, minimizing intermediate files.
        """
        import ffmpeg
        import os
        from pydub import AudioSegment
        temp_videos = []
        if overlay_texts is None:
            overlay_texts = [None] * len(image_files)
        audio_durations = [AudioSegment.from_file(audio).duration_seconds for audio in audio_files]
        # Generate temp videos for each slide
        for idx, (img, audio) in enumerate(zip(image_files, audio_files)):
            temp_video = os.path.join(os.path.dirname(output_path), f'temp_slide_{idx+1}.mp4')
            img_stream = ffmpeg.input(img, loop=1, framerate=1, t=audio_durations[idx] + slide_pause_duration)
            slide_audio = AudioSegment.from_file(audio)
            pause = AudioSegment.silent(duration=slide_pause_duration * 1000)
            slide_audio_with_pause = slide_audio + pause
            temp_audio_path = os.path.join(os.path.dirname(output_path), f'temp_audio_{idx+1}.mp3')
            slide_audio_with_pause.export(temp_audio_path, format='mp3')
            # Mix background music if provided
            if background_music_path and os.path.exists(background_music_path):
                music = AudioSegment.from_file(background_music_path)
                # Lower music volume and overlay
                music = music - 18
                mixed = slide_audio_with_pause.overlay(music)
                mixed.export(temp_audio_path, format='mp3')
            audio_stream = ffmpeg.input(temp_audio_path)
            (
                ffmpeg
                .output(img_stream, audio_stream, temp_video, vcodec='libx264', acodec='aac', pix_fmt='yuv420p', shortest=None, r=30)
                .overwrite_output()
                .run()
            )
            temp_videos.append(temp_video)
            os.remove(temp_audio_path)
        # Batch concat and fade transitions using ffmpeg complex filter
        # Build filter_complex for fade transitions
        filter_complex = ''
        inputs = []
        for i, v in enumerate(temp_videos):
            inputs.append(f'-i "{v}"')
        for i in range(len(temp_videos)):
            filter_complex += f'[{i}:v][{i}:a]'
        filter_complex += f'concat=n={len(temp_videos)}:v=1:a=1[v][a]'
        ffmpeg_cmd = f'ffmpeg -loglevel error ' + ' '.join(inputs) + f' -filter_complex "{filter_complex}" -map "[v]" -map "[a]" -y "{output_path}"'
        # Run ffmpeg command
        os.system(ffmpeg_cmd)
        # Cleanup temp videos
        for v in temp_videos:
            os.remove(v)
