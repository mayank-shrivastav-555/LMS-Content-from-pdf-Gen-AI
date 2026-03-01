def generate_ppt(slides, ppt_output, background_image_path=None, doc_file=None):
    """
    Wrapper to create a PPT using PPTCreator and save it to ppt_output.
    slides: list of dicts with slide content
    ppt_output: output pptx file path
    background_image_path: optional background image
    doc_file: optional document file for reference
    Returns: slide_chunks (list of bullet points per slide)
    """
    ppt = PPTCreator(background_image_path)
    # Add title slide
    if slides and 'title' in slides[0]:
        ppt.add_title_slide(slides[0]['title'], slides[0].get('subtitle', ''))
    heading_audio_list = []
    all_bullets = []
    for slide in slides:
        bullets = slide.get('bullets', [])
        if bullets:
            chunks = ppt.add_bullet_slide(bullets, title=slide.get('title'), subtitle=slide.get('subtitle'), heading_audio_list=heading_audio_list)
            all_bullets.extend(chunks)
    ppt.add_thankyou_slide()
    ppt.save(ppt_output)
    return all_bullets
from pptx import Presentation
from pptx.util import Inches, Pt


import re
from typing import List, Optional

class PPTCreator:
    def add_thankyou_slide(self):
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        self._set_slide_background(slide, self.CORPORATE_WHITE)
        # Calculate box size for central alignment
        box_width = Pt(600)
        box_height = Pt(120)
        left = int((self.slide_width - box_width) / 2)
        top = int((self.slide_height - box_height) / 2)
        txBox = slide.shapes.add_textbox(left, top, box_width, box_height)
        tf = txBox.text_frame
        tf.text = "Thank You!"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(60)
        run.font.bold = True
        run.font.color.rgb = self._rgb(*self.CORPORATE_AQUA)
    """
    Professional PowerPoint creator for dynamic, well-formatted presentations.
    """
    CORPORATE_BLUE = (0, 51, 102)
    CORPORATE_ORANGE = (255, 140, 0)
    CORPORATE_YELLOW = (255, 180, 40)
    CORPORATE_GRAY = (240, 240, 240)
    CORPORATE_WHITE = (255, 255, 255)
    CORPORATE_AQUA = (103, 255, 242)
    TITLE_FONT_SIZE = Pt(40)
    BULLET_FONT_SIZE = Pt(18)

    def __init__(self, background_image_path=None):
        self.prs = Presentation()
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        self.background_image_path = background_image_path

    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        self._set_slide_background(slide, self.CORPORATE_BLUE)
        # Use a more meaningful default title if 'Title' is passed
        slide_title = title
        if title.strip().lower() == "title":
            slide_title = "Policy Overview"
        self._set_text(slide.shapes.title, slide_title, self.TITLE_FONT_SIZE, self._rgb(*self.CORPORATE_AQUA), bold=True)
        self._set_text(slide.placeholders[1], subtitle, Pt(24), self._rgb(*self.CORPORATE_WHITE))


    def add_bullet_slide(self, bullet_points: list, title: str = None, subtitle: str = None, max_points_per_slide: int = 6, heading_audio_list: list = None) -> list:
        """
        Adds bullet slides and returns a list of chunks (list of bullet points per slide) for audio sync.
        Supports custom title and subtitle for the slide.
        """
        cleaned_points, slide_titles = self._process_bullet_points(bullet_points)
        slide_chunks = []
        for i in range(0, len(cleaned_points), max_points_per_slide):
            chunk = cleaned_points[i:i+max_points_per_slide]
            slide_chunks.append(chunk)
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
            slide_title = title if title else self._get_slide_title(i, max_points_per_slide, slide_titles, chunk)
            slide_subtitle = subtitle if subtitle else ""
            self._set_slide_background(slide, self.CORPORATE_WHITE)
            from pptx.util import Pt
            heading_font_size = Pt(32)            # shift the slide title down a bit to avoid overlap with background/logo
            try:
                title_shape = slide.shapes.title
                title_shape.top = title_shape.top + Inches(2.0)
            except Exception:
                pass
            self._set_text(slide.shapes.title, slide_title, heading_font_size, self._rgb(*self.CORPORATE_AQUA), bold=True)
            bullet_placeholder = slide.shapes.placeholders[1]
            # Increase spacing between heading and bullet content (e.g., 1.5 inch)
            from pptx.util import Inches
            bullet_placeholder.top = slide.shapes.title.top + slide.shapes.title.height + Inches(1.5)
            bullet_placeholder.left = slide.shapes.placeholders[1].left
            bullet_placeholder.width = slide.shapes.placeholders[1].width
            bullet_placeholder.height = slide.shapes.placeholders[1].height
            self.BULLET_FONT_SIZE = Pt(22)
            self._add_bullets(bullet_placeholder, chunk)
            # Optionally set subtitle if provided
            if subtitle:
                try:
                    self._set_text(slide.placeholders[1], slide_subtitle, self.BULLET_FONT_SIZE, self._rgb(120, 120, 120))
                except Exception:
                    pass
            # Add heading to heading_audio_list for narration
            if heading_audio_list is not None:
                heading_audio_list.append(slide_title)
        return slide_chunks

    def save(self, file_path: str) -> None:
        self.prs.save(file_path)

    @staticmethod
    def _rgb(r: int, g: int, b: int):
        from pptx.dml.color import RGBColor
        return RGBColor(r, g, b)

    def _set_slide_background(self, slide, color_rgb):
        # Set background color first
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(*color_rgb)
        # If you have a background image, add it behind all content
        if self.background_image_path:
            left = top = Inches(0)
            width = self.slide_width
            height = self.slide_height
            try:
                pic = slide.shapes.add_picture(self.background_image_path, left, top, width=width, height=height)
                # Move image to back so it does not cover text
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)
                pic.fill.transparency = 0.15
            except Exception:
                pass

    def _set_text(self, shape, text: str, font_size: Pt, color, bold: bool = False) -> None:
        shape.text = text
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.color.rgb = color
            paragraph.font.bold = bold

    def _process_bullet_points(self, bullet_points: List[str]) -> tuple[List[str], List[str]]:
        cleaned_points = []
        slide_titles = []
        for point in bullet_points:
            clean = re.sub(r'[\*#]+', '', point).strip()
            if (clean.lower().startswith('of course') or clean.startswith('---') or
                clean.lower().startswith('presentation:') or
                clean.lower().startswith('here is the summary')):
                continue
            # Filter out slide titles like 'Slide 1', 'Slide 2', etc.
            m = re.match(r"Slide \d+: (.+)", clean)
            if m:
                possible_title = m.group(1).strip()
                # Only add if not generic (not just a number or 'Slide X')
                if not re.match(r"^Slide \d+$", possible_title, re.IGNORECASE):
                    slide_titles.append(possible_title)
                continue
            h = re.match(r"\((Title|Subtitle)\): (.+)", clean)
            if h:
                possible_title = h.group(2).strip()
                if not re.match(r"^Slide \d+$", possible_title, re.IGNORECASE):
                    slide_titles.append(possible_title)
                continue
            # Also skip points that are just 'Slide X' as content
            if re.match(r"^Slide \d+$", clean, re.IGNORECASE):
                continue
            cleaned_points.append(clean)
        return cleaned_points, slide_titles

    def _get_slide_title(self, idx: int, max_points: int, slide_titles: List[str], chunk: List[str]) -> str:
        title_idx = idx // max_points
        if title_idx < len(slide_titles):
            return slide_titles[title_idx]
        # Use first heading-like bullet as title
        for point in chunk:
            if re.match(r'^[A-Z][^:]+:', point) or re.search(r'\bPolicy\b', point):
                return point.split(':')[0] if ':' in point else point
        return "Summary"

    def _add_bullets(self, placeholder, bullet_points):
        from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
        tf = placeholder.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        tf.margin_left = Inches(0.5)
        tf.margin_right = Inches(0.3)
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Vertically center the bullet points
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = self.BULLET_FONT_SIZE
            p.font.color.rgb = self._rgb(*self.CORPORATE_WHITE)
            p.alignment = PP_ALIGN.LEFT  # Horizontally left align the bullet points
            
    def add_intro_image_slide(self, image_path: str) -> None:
        """Add a full-bleed slide showing image_path as the first/intro slide."""
        print(f"Adding intro image slide with image: {image_path}")
        try:
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank layout
            # adding picture accepts EMU lengths (prs.slide_width/height are EMU), use them to fill the slide
            slide.shapes.add_picture(image_path, 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
        except Exception:
            # fallback: try adding without explicit sizing
            try:
                slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
                slide.shapes.add_picture(image_path, 0, 0)
            except Exception:
                # silently ignore if the image can't be added
                pass